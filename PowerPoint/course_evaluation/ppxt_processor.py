import utils
import random
from colorama import Fore
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import ppxt_vision

def get_formatting(ppxt_filepath, output_file):
    df = pd.DataFrame(columns=['Font', 'Size', 'Colour', 'Bold', 'Italics', 'Underline'])
    ppxt_file = Presentation(ppxt_filepath)

    print(f"{Fore.GREEN}---------------------Checking Formatting---------------------{Fore.RESET}")
    for slide in ppxt_file.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        df = df.append({'Font': run.font.name, 'Size': run.font.size, 'Colour': run.font.color.rgb, 'Bold': run.font.bold, 'Italics': run.font.italic, 'Underline': run.font.underline}, ignore_index=True)
                        print(f"Font: {run.font.name}, 'Size': {run.font.size}, 'Colour': {run.font.color.rgb}, Bold: {run.font.bold}, Italic: {run.font.italic}, Underline: {run.font.underline}")
    return df

def process(ppxt_filepath):
    #this function processes a PowerPoint file and returns a dataframe with the slide number, title, slide text, and notes text

    presentation = Presentation(ppxt_filepath)

    #create an empty dataframe
    df = pd.DataFrame(columns=['Slide Number', 'Slide Type', 'Title', 'Subtitle', 'Slide Text', 'Notes Text'])
    slides_as_images_path = "C:\\Users\kevin\\Documents\\Coding\\Scripts\\PowerPoint\\course_evaluation\\ppxt_images"
    
    print(f"{Fore.GREEN}---------------------Begin Processing Slides---------------------{Fore.RESET}")

    for i, slide in enumerate(presentation.slides, start=1):

        smart_art_detected = detect_smart_art(slide, i) 

        slide_number = i

        slide_type = "No Type"
        title = "No Title"
        subtitle = "No Subtitle"
        slide_text = ""
        notes_text = ""

        for shape in slide.shapes:

            if shape.is_placeholder and shape.placeholder_format.idx == 0:  # idx 0 is the title placeholder
                title = shape.text

            #if there is a text box, save the text to subtitle
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                if "https://www" not in shape.text: #verify it's not an image reference
                    if shape.text.strip() != subtitle.strip():
                        subtitle = shape.text

            # Accessing the text within the slide
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = "".join(run.text for run in paragraph.runs)
                    if paragraph_text.strip() and paragraph_text != title:  # Check if the paragraph text is not empty
                        if paragraph_text.strip() != subtitle.strip():
                                #check if the text contains a link
                                print(paragraph_text)
                                if "http" in paragraph_text:
                                    split_paragraph = paragraph_text.split("http")
                                    post_http = split_paragraph[1].split(" ")
                                    paragraph_text = split_paragraph[0]
                                    for i in range(1, len(post_http)):
                                        paragraph_text += post_http[i] + " "

                                slide_text += paragraph_text +"\n"
            

            # Accessing the notes within the slide
            if notes_text == "" and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for paragraph in notes_slide.notes_text_frame.paragraphs:
                    if "http" in paragraph.text:
                        split_paragraph = paragraph.text.split("http")
                        post_http = split_paragraph[1].split(" ")
                        paragraph.text = split_paragraph[0]
                        for i in range(1, len(post_http)):
                            paragraph.text += post_http[i] + " "

                    notes_text += paragraph.text
                

            if "agenda" in title.lower():
                slide_type = "Agenda"
            elif title == "No Title" and subtitle == "Recap":
                slide_type = "Recap"
            elif title == "No Title" and "Legal Disclaimers" in slide_text:
                slide_type = "Legal"
            elif "discussion" in title.lower():
                slide_type = "Discussion"
            elif "activity" in title.lower():
                slide_type = "Activity"
            elif (slide_number == 1 and slide_text == "AI for Workforce") or (subtitle == "No Subtitle" and slide_text == ""):
                slide_type = "Title"
            elif (title == "Not Title" and slide_text == "") or (subtitle == "No Subtitle" and slide_text != ""):
                slide_type = "Transition"
            else:
                slide_type = "Content"
        #add slide 
        df = df.append({'Slide Number': slide_number, 'Slide Type': slide_type, 'Title': title, 'Subtitle': subtitle, 'Slide Text': slide_text, 'Notes Text': notes_text, 'Smart Art Detected': smart_art_detected}, ignore_index=True)
        print(f"{Fore.GREEN}---------------------Processing Complete---------------------{Fore.RESET}")
        utils.print_loader(len(df), i)

    return df


def detect_smart_art(slide, i):
    #this function detects if a slide contains smart art
    smart_art_tags = ["</p:graphicFrame>", "<p:grpSp>"]
    xml = slide._element.xml
    smart_art_detected = False
    for i in range (len(smart_art_tags)):
        if smart_art_tags[i] in xml:
            smart_art_detected = True
    return smart_art_detected