import glob
import os
import random
from pptx import Presentation
from pptx.oxml.ns import nsdecls, _nsmap

def add_alt_text_to_image(ppxt_file, slide_index, alt_text):

    # Add the adec namespace
    _nsmap["adec"] = "http://schemas.microsoft.com/office/drawing/2017/decorative"

    # Get the specified slide
    slide = presentation.slides[slide_index]

    # Iterate through all shapes in the slide
    for shape in slide.shapes:
        print("slide_index: ", slide_index)
        print("shape.shape_type: ", shape.shape_type)
        if shape.shape_type == 13:  # Check if the shape is an image
            #save the image to the local directory
            image = shape.image
            image_bytes = image.blob
            image_filename = f"alt_text_img/image_{slide_index}.png"
            #check if the file name already exists
            if os.path.exists(image_filename):
                image_filename = f"alt_text_img/image_{slide_index}_{random.randint(0,1000)}.png"
            with open(image_filename, "wb") as f:
                f.write(image_bytes)

            #check if an image is marked as decorative
            adec_decoratives = shape._element._nvXxPr.cNvPr.xpath(".//adec:decorative[@val='1']")
            if not adec_decoratives:
                if shape._element._nvXxPr.cNvPr.attrib.get("descr", "") == "":
                    shape._element._nvXxPr.cNvPr.attrib["descr"] = alt_text  # Set the alt text for the image
            
    #run through the folder alt_text_img and delete all the images
    img_files = glob.glob(os.path.join("alt_text_img/", "*"))
    for img_file in img_files:
        os.remove(img_file)
    

    # Save the modified presentation
    presentation.save(ppxt_file)

# Usage example
ppxt_file = "C:/Users/kevin/Downloads/Week 06.pptx"
presentation = Presentation(ppxt_file)
alt_text = '2nd Alt Text Run'

for i in range(len(presentation.slides)):
    slide_i = i
    add_alt_text_to_image(ppxt_file, slide_i, alt_text)