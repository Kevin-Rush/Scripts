import pptx
from openai import OpenAI

import os
from dotenv import load_dotenv


load_dotenv()
gpt_api_key = os.getenv("FS_LLM_EVAL_OPENAI_API_KEY")

client = OpenAI()
client.api_key = gpt_api_key
# Read the PowerPoint file
presentation = pptx.Presentation(r"C:\Users\kevin\Downloads\Artificial Intelligence (AI) For DevOps - Trainer PPT.pptx")

# Iterate through each slide
for slide in presentation.slides:
    # Read the content on the slide
    slide_content = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_content += run.text

    # Generate script using OpenAI API
    prompt_tokens = 0
    completion_tokens = 0
    model = "gpt-4-0125-preview"
    messages = [
        {"role": "system", "content": "You are an expert script writing assistant that excels at interpreting presenter slides to create cohesive and elegant speaker notes for presenters to read."},
        {"role": "user", "content": "Can you please use this slide content to generate the speaker's note to go along with it: " + slide_content}
    ]

    response = client.chat.completions.create(
        model=model,
        messages=messages,
    )
    #messages.pop()
    prompt_tokens += response.usage.prompt_tokens
    completion_tokens += response.usage.completion_tokens
    response = response.choices[0].message.content

    # Write the generated script into the slide notes
    slide.notes_slide.notes_text_frame.text = response

# Save the modified PowerPoint file
presentation.save(r"C:\Users\kevin\Downloads\Artificial Intelligence (AI) For DevOps - Trainer PPT.pptx")

