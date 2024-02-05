from pptx import Presentation
import openai
import pandas as pd

def evaluate_ppxt(ppxt_filepath):
    presentation = Presentation(ppxt_filepath)
    #read in the api key from the file gpt_api_key.txt in the folder above the current folder
    with open("C:/Users/kevin/Documents/Coding/Scripts/gpt_api_key.txt", "r") as file:
        openai.api_key = file.read()


    #create an empty dataframe
    df = pd.DataFrame(columns=['Slide Number', 'Title', 'Slide Text', 'Notes Text', 'Simplicity', 'Clarity', 'Overall'])

    for i, slide in enumerate(presentation.slides, start=1):
        slide_number = i
        title = "No Title"
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:  # idx 0 is the title placeholder
                title = shape.text
                break

        # Accessing text within the slide
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = "".join(run.text for run in paragraph.runs)
                    if paragraph_text.strip():  # Check if the paragraph text is not empty
                        slide_text += "* " + paragraph_text + "\n"
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_text += "* " + run.text + "\n"

        # Accessing notes within the slide
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for paragraph in notes_slide.notes_text_frame.paragraphs:
                notes_text += paragraph.text


        if title == "Agenda":
            prompt_agenda = "Hello ChatGPT, I am a powerpoint evaluator and I'm reviewing a new course to ensure it is well written. I'm going to pass you an agenda slide and I'd like you to tell me if the flow of this section is logical and the section names are informative and direct. For example, a common issue I see is fanciful naming like: Mastering the Art and Science of Prompt Engineering. This is a extravegant title that does not actually convey what is going to be covered in this section. I want the agenda slide to be informative and direct. Please only consider the text I've given, and respond in a tone that is appropriate for giving feedback to a colleague. \n\n Slide Contents: " + slide_text
            
            response = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=[{"role": "user", "content": prompt_simplicity}])
            response = response['choices'][0]['message']['content']
        else:
            # Evaluate the slide using chatgpt
            prompt_simplicity = "Hello ChatGPT, I am a powerpoint evaluator and I'm reviewing a new course to ensure it is well written. I'm going to pass you the slide title and text and I'd like you to tell me if the language used is the simplest and mosth pithy way of conveying the information on the slide. For example, a common issue I see is the overuse of unnecessary but formal sounding language like the word utilize. Although utilize is an accurate word, in most cases use is a better word because it is simple, short, and commmon. Please be on the look out for unnecessary language. I want the slides to be professional, pithy, and direct. Please only consider the text I've given, and respond in a tone that is appropriate for giving feedback to a colleague. \n\n Title:" + title + " \nSlide Contents: " + slide_text
            
            response = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=[{"role": "user", "content": prompt_simplicity}])
            response = response['choices'][0]['message']['content']

        #add slide 
        df = df.append({'Slide Number': slide_number, 'Title': title, 'Slide Text': slide_text, 'Notes Text': notes_text, 'Feedback':response}, ignore_index=True)

        # Print the slide information
        print(f"Slide Number: {slide_number}")
        print(f"Title: {title}")
        print(f"Slide Text: {slide_text}")
        print(f"Notes Text: {notes_text}")
        print(f"Feedback: {response}")
        print()        
        

# Usage
ppxt_file = "C:/Users/kevin/Downloads/test_file.pptx"
evaluate_ppxt(ppxt_file)