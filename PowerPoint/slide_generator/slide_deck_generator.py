from pptx import Presentation
import json

def generate_presentation(json_script, base_ppxt):
    #Create a presentation object
    presentation = Presentation(base_ppxt)

    #read in the script from the json file
    with open(json_script, 'r') as file:
        content = json.load(file)

    for i in range(len(content)):

        if content[i]['title'] == "DECK TITLE":
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            title = slide.shapes.title
            title.text = content[i]['content']

            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content[i]['notes']
            
        elif content[i]['title'] == "Agenda" or content[i]['title'] == "Discussion" or content[i]['title'] == "Review of Last Week" or content[i]['title'] == "Key Takeaways":
            slide = presentation.slides.add_slide(presentation.slide_layouts[3])
            title = slide.shapes.title
            title.text = content[i]['title']

            slide_content = slide.shapes[1]
            slide_content.text = content[i]['content']

            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content[i]['notes']
            
        elif content[i]['title'] == "TRANSITION":
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

            slide_content = slide.shapes[0]
            slide_content.text = content[i]['content']

            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content[i]['notes']

        elif content[i]['title'] == "RECAP":
            slide = presentation.slides.add_slide(presentation.slide_layouts[2])
        
        elif content[i]['title'] == "LEGAL":
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])

        elif content[i]['title'] == "Activity":
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            title = slide.shapes.title
            title.text = content[i]['title']

            slide_subtitle = slide.shapes[1]
            slide_subtitle.text = content[i]['subtitle']

            slide_content = slide.shapes[2]
            slide_content.text = content[i]['content']

            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content[i]['notes']
        else:   
            slide = presentation.slides.add_slide(presentation.slide_layouts[4])
            title = slide.shapes.title
            title.text = content[i]['title']

            slide_subtitle = slide.shapes[1]
            slide_subtitle.text = content[i]['subtitle']

            slide_content = slide.shapes[2]
            slide_content.text = content[i]['content']

            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = content[i]['notes']

    # Save the presentation
    presentation.save("generated_presentation.pptx")