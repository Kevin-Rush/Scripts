from pptx import Presentation
import pandas as pd

def process_ppxt(ppxt_filepath):
    presentation = Presentation(ppxt_filepath)

    #create an empty dataframe
    df = pd.DataFrame(columns=['Slide Number', 'Title', 'Slide Text', 'Notes Text'])

    for i, slide in enumerate(presentation.slides, start=1):
        slide_number = i
        title = "No Title"
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:  # idx 0 is the title placeholder
                title = shape.text
                break

        if "agenda" in title.lower() or "discussion" in title.lower() or "activity" in title.lower():
            
            # Accessing text within the slide
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = "".join(run.text for run in paragraph.runs)
                        if paragraph_text.strip():  # Check if the paragraph text is not empty
                            slide_text += paragraph_text + "\n"
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    slide_text += run.text + "\n"

            # Accessing notes within the slide
            notes_text = ""
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                for paragraph in notes_slide.notes_text_frame.paragraphs:
                    notes_text += paragraph.text


            #add slide 
            df = df.append({'Slide Number': slide_number, 'Title': title, 'Slide Text': slide_text, 'Notes Text': notes_text}, ignore_index=True)

    return df     
        

def run(ppxt_file):
    return evaluate_ppxt(ppxt_file)
