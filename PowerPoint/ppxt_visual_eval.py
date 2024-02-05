from pptx import Presentation
from PIL import Image
import io

def save_slides_as_images(pptx_file):
    presentation = Presentation(pptx_file)
    
    for i, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_file = io.BytesIO(image.blob)
                pil_image = Image.open(image_file)
                pil_image.save(f"slide_{i+1}.png")

# Usage
ppxt_file = "C:/Users/kevin/Downloads/test_file.pptx"
save_slides_as_images(ppxt_file)