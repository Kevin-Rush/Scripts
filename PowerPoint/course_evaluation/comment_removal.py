import pptx
import aspose.slides
# from aspose.words import Document, NodeType


def extract_comments_from_pptx(pptx_file, output_file):
    # Load the PowerPoint presentation
    presentation = aspose.slides.Presentation(pptx_file)

    # Open the output file in write mode
    with open(output_file, 'w') as file:

        # comments = doc.get_child_nodes(NodeType.COMMENT, True)

        # Loop through the list of authors
        for author in presentation.comment_authors:
            # Loop through each author's comments
            for comment in author.comments:
                # Access comment properties
                print(f"Comment by {author.name}, at {comment.created_time}: {comment.text}")
                # Write the comment to the output file
                file.write(f"Comment by {author.name}, at {comment.created_time}: {comment.text}\n")


            author.comments.clear()
    presentation.comment_authors.clear()

    # Save the modified presentation
    presentation.save(pptx_file, aspose.slides.export.SaveFormat.PPTX)
    remove_aspose_watermark(pptx_file)

def remove_aspose_watermark(pptx_file):
    # Accessing the text within the slide
    presentation = pptx.Presentation(pptx_file)
    found = False
    for slide in presentation.slides:
        for i in range(len(slide.shapes)):
            try:
                shape = slide.shapes[i]
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_text = "".join(run.text for run in paragraph.runs)
                        # Check if the paragraph text contains the Aspose watermark
                        if "Created with Aspose.Slides for Python" in paragraph_text:
                            print("Found Aspose watermark")
                            # Remove the Aspose watermark
                            found = True
                if found:
                    #clear the entire text frame
                    shape.text_frame.clear()
                    found = False
            except:
                print("Error accessing shape")
                continue
    
    # Save the modified presentation
    presentation.save(pptx_file)


# # Specify the input PowerPoint file and output text file
# pptx_file = "C:\\Users\\kevin\\Downloads\\comments_test.pptx"
# output_file = 'comments_output.txt'

# # Call the function to extract comments and remove them from the PowerPoint file
# extract_comments_from_pptx(pptx_file, output_file)