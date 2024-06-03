import requests
# from sapling import SaplingClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
from pptx import Presentation

def evaluate(ppxt_filepath, sapling_api_key):

    presentation = Presentation(ppxt_filepath)
    client = SaplingClient(api_key=sapling_api_key)

    #create a notes file
    changes_made = open(ppxt_filepath[:-5] + "_changes_made.txt", "w", encoding="utf-8")

    for i, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:

            #access the text in all types of shapes
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:                        
                        run_text = run.text
                        edits = client.edits(run_text, session_id='test_session')

                        if run_text == "Legal Disclaimers":
                            continue
                        elif edits["edits"] != []:
                            print(f'Slide: {i}')
                            print(f'run_text: {run_text}')
                            print(f'Edits: {edits["edits"]}')
                            for edit in edits["edits"]:
                                start = edit['sentence_start'] + edit['start']
                                end = edit['sentence_start'] + edit['end']
                                if start > len(run_text) or end > len(run_text):
                                    print(f'Edit start:{start}/end:{end} outside of bounds of text:{run_text}')
                                    continue
                                run_text = run_text[: start] + edit['replacement'] + run_text[end:]
                            
                            #add the slide number, the old text and the new text to the changes made file
                            changes_made.write(f"Slide: {i}\nOld Text: {run.text}\nNew Text: {run_text}\n\n")
                            run.text = run.text.replace(run.text, run_text)
    #save the edited presentation to a file with the same path but a slightly different name
    presentation.save(ppxt_filepath[:-5] + "_edited.pptx")
    changes_made.close()
    print("Grammar evaluation complete")

def check_font_type(ppxt_filepath):
    presentation = Presentation(ppxt_filepath)
    i = 1
    for slide in presentation.slides:
        if i > 4:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            print(run.font.name)
                            if shape.is_placeholder and shape.placeholder_format.idx == 0:  # idx 0 is the title placeholder
                                if run.font.name != "IntelOne Display Light":
                                    print("Wrong font detected!")
                                    print("Slide Number: ", i)
                                    print("Font: ", run.font.name)
                                    print("Text: ", run.text)

                                    #change the font to the correct font
                                    run.font.name = "IntelOne Display Light"
                                    print("Font changed to IntelOne Display Light")
                            else:
                                if run.font.name != "IntelOne Display Regular":
                                    print("Wrong font detected!")
                                    print("Slide Number: ", i)
                                    print("Font: ", run.font.name)
                                    print("Text: ", run.text)

                                    #change the font to the correct font
                                    run.font.name = "IntelOne Display Regular"
                                    print("Font changed to IntelOne Display Regular")
        i += 1
    presentation.save(ppxt_filepath)
