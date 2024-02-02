from pptx import Presentation

def extract_script(ppxt_file, output_file):
    presentation = Presentation(ppxt_file)
    script = ""

    for i, slide in enumerate(presentation.slides, start=1):
        # Find the title for the slide
        title = "No Title"
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:  # idx 0 is the title placeholder
                title = shape.text
                break

        # Add the slide number and title to the script
        script += f"\n\nSlide {i} - {title}\n"

        # Extract the speaker notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for paragraph in notes_slide.notes_text_frame.paragraphs:
                for run in paragraph.runs:
                    script += run.text + " "

    with open(output_file, "w", encoding='utf-8') as file:
        file.write(script)

ppxt_file = "C:/Users/kevin/Downloads/Intel AI for Workforce TTT AEAI.pptx"
output_file = "C:/Users/kevin/Documents/Coding/Scripts/PowerPoint/output_extract_script.txt"
extract_script(ppxt_file, output_file)