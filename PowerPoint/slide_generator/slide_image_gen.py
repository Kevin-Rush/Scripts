import json
from sys import displayhook
import openai
from openai import OpenAI
from pptx import Presentation
import requests
from PIL import Image
from pptx.util import Inches

def  gen_image_for_slide(json_script, presentation_file, api_key):

    client = OpenAI()
    client.api_key = api_key

    #Create a presentation object
    presentation = Presentation(presentation_file)

    #read in the script from the json file
    with open(json_script, 'r') as file:
        content = json.load(file)

    for i in range(len(content)):
        #load the slide number corresponding to i
        print("-----------------Slide: ", i+1, "-----------------")
        slide = presentation.slides[i]
        notes = content[i]['notes']
        print()
        print("Notes: ")
        print(notes)
        
        image_description = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an expert text for image interpreter. A text for image interpreter, takes text from all contexts (business, IT, science, biology, etc.) and is able to image an image that best represents the interpreted text."}, 
                {"role": "user", "content": "Hello, can you please expertly interpret my text and tell me what image best represents the text? Here is my text: "+notes}],
        )

        image_description = image_description.choices[0].message.content
        print("-----------------Image Description-----------------")
        print(image_description)

        response = client.images.generate(
            model="dall-e-3",
            prompt=image_description,
            size="1024x1024",
            quality="standard",
            n=1,
        )

        # save the image
        generated_image_name = f"slide_{i}.png"  
        generated_image_url = response.data[0].url  # extract image URL from response
        generated_image = requests.get(generated_image_url).content  # download the image

        with open(generated_image_name, "wb") as image_file:
            image_file.write(generated_image)  # write the image to the file
        
        #resize the image to 256x256
        img = Image.open(generated_image_name)
        img = img.resize((256, 256))
        img.save(generated_image_name)

        #add the image to the presentation object
        slide = presentation.slides[i]
        left = Inches(7.48)
        top = Inches(2.33)
        pic = slide.shapes.add_picture(generated_image_name, left, top)
        
        presentation.save(presentation_file)