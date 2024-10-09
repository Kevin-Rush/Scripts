import os
from pptx import Presentation


from .api_handling import make_get_request
from . import party_functions as party_fns
from . import org_functions as org_fns

def paginate_get_request(endpoint):
    """
    Handles paginated GET requests from the Capsule API.
    
    :param endpoint: The API endpoint to target.
    :return: List of all results across pages.
    """
    results = []
    page = 1
    more_pages = True

    while more_pages:
        data = make_get_request(f'{endpoint}?page={page}')
        if data and 'entries' in data:
            results.extend(data['entries'])
            if len(data['entries']) < 50:
                more_pages = False
            else:
                page += 1
        else:
            more_pages = False

    return results

def format_attachments(attachments):
    """
    Formats the attachments information into a readable string.

    :param attachments: List of attachment dictionaries.
    :return: A formatted string for attachments.
    """
    if attachments:
        return ', '.join([attachment['filename'] for attachment in attachments])
    return "No attachments"


def save_professor_ids_to_file(filename="professors_ids.txt"):
    """
    Fetches all professors and saves their names and IDs to a file.
    
    :param filename: The name of the file to save professor details to.
    """
    professors = party_fns.get_all_people()  # Fetches all professors
    if professors:
        with open(filename, 'w', encoding='utf-8') as f:
            f.write("Professor Name, Professor ID\n")
            f.write("=" * 40 + "\n")
            for professor in professors:
                professor_name = f"{professor.get('firstName', 'Unknown')} {professor.get('lastName', 'Unknown')}"
                professor_id = professor.get('id')
                f.write(f"{professor_name}, {professor_id}\n")
        print(f"Saved professor names and IDs to {filename}")
    else:
        print("No professors found.")


def save_college_ids_to_file(filename="college_ids.txt"):
    """
    Fetches all colleges and saves their names and IDs to a file.
    
    :param filename: The name of the file to save college details to.
    """
    colleges = org_fns.get_all_organizations()  # Fetches all colleges
    if colleges:
        with open(filename, 'w', encoding='utf-8') as f:
            f.write("College Name, College ID\n")
            f.write("=" * 40 + "\n")
            for college in colleges:
                college_name = college.get('name', 'Unknown College')
                college_id = college.get('id')
                f.write(f"{college_name}, {college_id}\n")
        print(f"Saved college names and IDs to {filename}")
    else:
        print("No colleges found.")

def extract_pptx_data(pptx_file):
    prs = Presentation(pptx_file)
    slide_data = {}

    for slide in prs.slides:
        title = ''
        body_texts = []

        # Attempt to get the slide title
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        else:
            # If there's no title placeholder, use the first shape with text as the title
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    title = shape.text.strip()
                    break

        # Collect body text from all text shapes, excluding the title
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                if shape != slide.shapes.title:
                    body_texts.append(shape.text.strip())

        body_text = '\n'.join(body_texts).strip()
        slide_data[title] = body_text

    return slide_data




pptx_file = r'C:\Users\kevin\Documents\Coding\Scripts\crm\search\log_presentation_working.pptx'
data = extract_pptx_data(pptx_file)
for college_name, slide_text in data.items():
    print(f"College Name: {college_name}")
    print(f"Slide Text: {slide_text}\n")

