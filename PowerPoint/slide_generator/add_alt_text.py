import base64
from pptx import Presentation
from colorama import Fore
import openai

# Function to encode the image
def encode_image(image_path):
    #This function encodes the image to base64
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def add_alt_text_to_images(pptx_file, api_key):
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file)

    openai.api_key = api_key

    # Iterate through each slide in the presentation
    for slide in presentation.slides:
        # Iterate through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape is an image
            if shape.shape_type == 13:  # 13 represents an image shape
                # Check if the image has alt text
                if not shape.has_alt_text:
                    # Add alt text to the image

                    # Get the image from shape
                    image = shape.image
                    base64_image = encode_image(image)

                    headers = {
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {api_key}"
                    }

                    payload = {
                    "model": "gpt-4-vision-preview",
                    "messages": [
                        {
                        "role": "user",
                        "content": [
                            {
                            "type": "text",
                            "text": "You are an alt text assistant. Your job is to review an image from an educational slide deck and describe the image best for visually impared students. You are an expert at being direct, pithy, and capturing the essence of an image in the shortest sentence possible. Can you please describe the image in the most direct and pithy way possible?"
                            },
                            {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}"
                            }
                            }
                        ]
                        }
                    ],
                    "max_tokens": 1251
                    }

                    response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
                    response = response.json()
                    
                    shape.alt_text = response

    # Save the modified presentation
    presentation.save("modified_presentation.pptx")


# Usage
with open("C:/Users/kevin/Documents/Coding/Scripts/gpt_api_key.txt", "r") as file:
    gpt_api_key = file.read()

add_alt_text_to_images("C:\Users\kevin\Downloads\alt_text_test.pptx", gpt_api_key)
